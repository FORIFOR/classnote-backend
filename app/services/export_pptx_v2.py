"""export_pptx_v2.py — Presentation-IR driven PPTX renderer.

Takes the dict produced by `presentation_ir.build_presentation_ir` and draws
16:9 slides in an "executive-dark" theme. Each slide kind has a dedicated
renderer so layouts stay centralised and visually consistent.

This module intentionally does NOT reinterpret summaryJson semantics; all
section-level decisions live in presentation_ir.py.

Supported kinds (Phase 1):
  - hero_summary
  - decision_board
  - todo_board
  - risk_alert
  - keyword_cloud
  - appendix_transcript
"""

from __future__ import annotations

import io
import logging
from typing import Any, Dict, List, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Theme: executive-dark
# ---------------------------------------------------------------------------

C_BG          = RGBColor(0x0F, 0x11, 0x15)
C_SURFACE     = RGBColor(0x17, 0x1A, 0x21)
C_SURFACE_ALT = RGBColor(0x1D, 0x22, 0x30)
C_TEXT        = RGBColor(0xF5, 0xF7, 0xFB)
C_TEXT_MUTED  = RGBColor(0xA7, 0xB0, 0xC0)
C_PRIMARY     = RGBColor(0x6E, 0xA8, 0xFE)
C_SUCCESS     = RGBColor(0x34, 0xC7, 0x59)
C_WARNING     = RGBColor(0xFF, 0xB0, 0x20)
C_DANGER      = RGBColor(0xFF, 0x5D, 0x5D)
C_BORDER      = RGBColor(0x2B, 0x32, 0x42)
C_CHIP_BG     = RGBColor(0x24, 0x2B, 0x3C)

FONT_HEADING = "Yu Gothic"
FONT_BODY    = "Yu Gothic"

# Slide geometry (16:9 default from python-pptx: 9144000 x 6858000 EMU ≈ 10" x 5.63")
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
MARGIN_X = Inches(0.55)
MARGIN_TOP = Inches(0.45)

STATUS_COLORS = {
    "confirmed": C_SUCCESS,
    "open": C_WARNING,
    "rejected": C_DANGER,
    "inferred": C_TEXT_MUTED,
}

SEVERITY_COLORS = {
    "high": C_DANGER,
    "medium": C_WARNING,
    "low": C_TEXT_MUTED,
}

PRIORITY_COLORS = {
    "high": C_DANGER,
    "normal": C_PRIMARY,
    "low": C_TEXT_MUTED,
}


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def render_pptx_from_ir(ir: Dict[str, Any]) -> bytes:
    """Render a Presentation IR into PPTX bytes (executive-dark theme).

    2026-05-08 design refresh: the deck is capped at **1–2 slides**
    (cover + a single ``hero_summary``). PPTX is now treated as a
    summary attachment; per-section content (decisions / todos /
    risks / keyword cloud / transcript appendix) lives in the PDF /
    DOCX exports. The per-kind handlers are kept in this module so a
    future PR can re-enable richer decks via a feature flag without
    re-writing the renderers.
    """
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    meta = ir.get("meta") or {}
    deck = ir.get("deck") or {}
    slides = ir.get("slides") or []

    # Cover slide always present
    _render_cover(prs, deck=deck, meta=meta)

    # Pick at most one ``hero_summary`` slide. Falls back to the first
    # slide of any recognised kind so we always have *something* on
    # slide 2 — but only ever one.
    hero = next((s for s in slides if s.get("kind") == "hero_summary"), None)
    if hero is None and slides:
        hero = slides[0]
    if hero is not None:
        try:
            kind = hero.get("kind")
            if kind == "hero_summary":
                _render_hero_summary(prs, hero)
            elif kind == "decision_board":
                _render_decision_board(prs, hero)
            elif kind == "todo_board":
                _render_todo_board(prs, hero)
            elif kind == "risk_alert":
                _render_risk_alert(prs, hero)
            elif kind == "keyword_cloud":
                _render_keyword_cloud(prs, hero)
            elif kind == "appendix_transcript":
                _render_appendix_transcript(prs, hero)
            else:
                _render_hero_summary(prs, hero)  # best-effort fallback
        except Exception as exc:
            logger.exception(f"[pptx_v2] failed to render summary slide: {exc}")

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Cover
# ---------------------------------------------------------------------------

def _render_cover(prs: Presentation, *, deck: Dict[str, Any], meta: Dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _paint_background(slide, C_BG)

    # Left accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.18), Inches(SLIDE_H_IN),
    )
    _fill(bar, C_PRIMARY)
    _no_line(bar)

    title = deck.get("title") or meta.get("title") or "Untitled"
    subtitle = deck.get("subtitle") or ""

    _add_textbox(
        slide,
        left=Inches(1.1), top=Inches(2.4),
        width=Inches(SLIDE_W_IN - 2.2), height=Inches(1.4),
        text=title,
        font_size=Pt(40), bold=True, color=C_TEXT, font=FONT_HEADING,
    )
    if subtitle:
        _add_textbox(
            slide,
            left=Inches(1.1), top=Inches(3.8),
            width=Inches(SLIDE_W_IN - 2.2), height=Inches(0.6),
            text=subtitle,
            font_size=Pt(18), color=C_TEXT_MUTED, font=FONT_HEADING,
        )

    chips: List[str] = []
    for key in ("dateText", "durationText"):
        v = meta.get(key)
        if v:
            chips.append(str(v))
    mode = meta.get("mode")
    if mode == "meeting":
        chips.append("会議")
    elif mode == "lecture":
        chips.append("講義")
    if chips:
        _draw_chip_row(slide, left=Inches(1.1), top=Inches(4.6),
                       width=Inches(SLIDE_W_IN - 2.2), chips=chips)

    _add_footer(slide, deck.get("footerText"))


# ---------------------------------------------------------------------------
# hero_summary
# ---------------------------------------------------------------------------

def _render_hero_summary(prs: Presentation, slide_ir: Dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _paint_background(slide, C_BG)
    _draw_slide_header(slide, slide_ir.get("title") or "結論")

    headline = slide_ir.get("headline") or ""
    subheadline = slide_ir.get("subheadline") or ""

    _add_textbox(
        slide,
        left=MARGIN_X, top=Inches(2.0),
        width=Inches(SLIDE_W_IN) - MARGIN_X * 2,
        height=Inches(2.4),
        text=headline,
        font_size=Pt(34), bold=True, color=C_TEXT, font=FONT_HEADING,
        word_wrap=True,
    )
    if subheadline:
        _add_textbox(
            slide,
            left=MARGIN_X, top=Inches(4.55),
            width=Inches(SLIDE_W_IN) - MARGIN_X * 2,
            height=Inches(1.5),
            text=subheadline,
            font_size=Pt(18), color=C_TEXT_MUTED, font=FONT_BODY,
            word_wrap=True,
        )

    chips = slide_ir.get("chips") or []
    if chips:
        _draw_chip_row(slide, left=MARGIN_X, top=Inches(6.3),
                       width=Inches(SLIDE_W_IN) - MARGIN_X * 2, chips=chips)

    _add_footer(slide)


# ---------------------------------------------------------------------------
# decision_board
# ---------------------------------------------------------------------------

def _render_decision_board(prs: Presentation, slide_ir: Dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _paint_background(slide, C_BG)
    _draw_slide_header(slide, slide_ir.get("title") or "決定事項")

    items = slide_ir.get("items") or []
    _draw_card_grid(slide, items, card_builder=_build_decision_card)

    _add_footer(slide)


def _build_decision_card(slide, x: Emu, y: Emu, w: Emu, h: Emu, item: Dict[str, Any]) -> None:
    _draw_card_base(slide, x, y, w, h)
    status = item.get("status") or "confirmed"
    accent = STATUS_COLORS.get(status, C_PRIMARY)

    _draw_card_accent_bar(slide, x, y, h, accent)

    label = item.get("label") or ""
    value = item.get("value") or ""
    reason = item.get("reason")
    owner = item.get("owner")
    due = item.get("due")

    pad = Inches(0.25)
    inner_x = x + Inches(0.15)
    inner_w = w - Inches(0.3)

    _add_textbox(
        slide,
        left=inner_x, top=y + pad,
        width=inner_w, height=Inches(0.35),
        text=label,
        font_size=Pt(11), color=C_TEXT_MUTED, font=FONT_BODY,
    )
    _add_textbox(
        slide,
        left=inner_x, top=y + Inches(0.55),
        width=inner_w, height=Inches(0.7),
        text=value,
        font_size=Pt(16), bold=True, color=C_TEXT, font=FONT_HEADING,
        word_wrap=True,
    )
    if reason:
        _add_textbox(
            slide,
            left=inner_x, top=y + Inches(1.25),
            width=inner_w, height=Inches(0.9),
            text=reason,
            font_size=Pt(10), color=C_TEXT_MUTED, font=FONT_BODY,
            word_wrap=True,
        )

    bottom_chips: List[Tuple[str, RGBColor]] = []
    if owner:
        bottom_chips.append((f"👤 {owner}", C_PRIMARY))
    if due:
        bottom_chips.append((f"📅 {due}", C_WARNING))
    if status:
        bottom_chips.append((_status_label(status), accent))
    if bottom_chips:
        _draw_colored_chips(slide, left=inner_x, top=y + h - Inches(0.55),
                            width=inner_w, chips=bottom_chips)


# ---------------------------------------------------------------------------
# todo_board
# ---------------------------------------------------------------------------

def _render_todo_board(prs: Presentation, slide_ir: Dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _paint_background(slide, C_BG)
    _draw_slide_header(slide, slide_ir.get("title") or "次のアクション")

    items = slide_ir.get("items") or []
    _draw_card_grid(slide, items, card_builder=_build_todo_card)

    _add_footer(slide)


def _build_todo_card(slide, x: Emu, y: Emu, w: Emu, h: Emu, item: Dict[str, Any]) -> None:
    _draw_card_base(slide, x, y, w, h)
    priority = item.get("priority") or "normal"
    accent = PRIORITY_COLORS.get(priority, C_PRIMARY)
    _draw_card_accent_bar(slide, x, y, h, accent)

    task = item.get("task") or ""
    owner = item.get("owner")
    due = item.get("due")
    status = item.get("status") or "open"
    blocking = bool(item.get("blocking"))

    inner_x = x + Inches(0.15)
    inner_w = w - Inches(0.3)

    _add_textbox(
        slide,
        left=inner_x, top=y + Inches(0.25),
        width=inner_w, height=Inches(1.2),
        text=task,
        font_size=Pt(15), bold=True, color=C_TEXT, font=FONT_HEADING,
        word_wrap=True,
    )

    meta_parts: List[Tuple[str, RGBColor]] = []
    if owner:
        meta_parts.append((f"👤 {owner}", C_PRIMARY))
    if due:
        meta_parts.append((f"📅 {due}", C_WARNING))
    pri_label = {"high": "優先度 高", "low": "優先度 低", "normal": "優先度 中"}.get(priority, priority)
    meta_parts.append((pri_label, accent))
    meta_parts.append((_status_label(status), STATUS_COLORS.get(status, C_TEXT_MUTED)))
    if blocking:
        meta_parts.append(("🚫 ブロッカー", C_DANGER))

    _draw_colored_chips(slide, left=inner_x, top=y + h - Inches(0.55),
                        width=inner_w, chips=meta_parts)


# ---------------------------------------------------------------------------
# risk_alert
# ---------------------------------------------------------------------------

def _render_risk_alert(prs: Presentation, slide_ir: Dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _paint_background(slide, C_BG)
    _draw_slide_header(slide, slide_ir.get("title") or "要確認・未決事項")

    items = slide_ir.get("items") or []
    _draw_card_grid(slide, items, card_builder=_build_risk_card)
    _add_footer(slide)


def _build_risk_card(slide, x: Emu, y: Emu, w: Emu, h: Emu, item: Dict[str, Any]) -> None:
    _draw_card_base(slide, x, y, w, h)
    severity = item.get("severity") or "medium"
    accent = SEVERITY_COLORS.get(severity, C_WARNING)
    _draw_card_accent_bar(slide, x, y, h, accent)

    title = item.get("title") or ""
    detail = item.get("detail")

    inner_x = x + Inches(0.15)
    inner_w = w - Inches(0.3)

    _add_textbox(
        slide,
        left=inner_x, top=y + Inches(0.2),
        width=inner_w, height=Inches(0.35),
        text=_severity_label(severity),
        font_size=Pt(11), bold=True, color=accent, font=FONT_BODY,
    )
    _add_textbox(
        slide,
        left=inner_x, top=y + Inches(0.55),
        width=inner_w, height=Inches(0.8),
        text=title,
        font_size=Pt(14), bold=True, color=C_TEXT, font=FONT_HEADING,
        word_wrap=True,
    )
    if detail:
        _add_textbox(
            slide,
            left=inner_x, top=y + Inches(1.35),
            width=inner_w, height=h - Inches(1.5),
            text=detail,
            font_size=Pt(10), color=C_TEXT_MUTED, font=FONT_BODY,
            word_wrap=True,
        )


# ---------------------------------------------------------------------------
# keyword_cloud (chip grid)
# ---------------------------------------------------------------------------

def _render_keyword_cloud(prs: Presentation, slide_ir: Dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _paint_background(slide, C_BG)
    _draw_slide_header(slide, slide_ir.get("title") or "キーワード")

    items = slide_ir.get("items") or []
    if not items:
        _add_footer(slide)
        return

    # Simple wrap-flow of chips
    x = MARGIN_X
    y = Inches(1.6)
    line_h = Inches(0.6)
    max_x = Inches(SLIDE_W_IN) - MARGIN_X

    for word in items:
        chip_w = _estimate_chip_width(word, Pt(14))
        if x + chip_w > max_x:
            x = MARGIN_X
            y += line_h
            if y + line_h > Inches(SLIDE_H_IN - 0.8):
                break
        _draw_chip(slide, x=x, y=y, width=chip_w, height=Inches(0.45),
                   text=word, fg=C_TEXT, bg=C_CHIP_BG, font_size=Pt(14))
        x += chip_w + Inches(0.12)

    _add_footer(slide)


# ---------------------------------------------------------------------------
# appendix_transcript
# ---------------------------------------------------------------------------

def _render_appendix_transcript(prs: Presentation, slide_ir: Dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _paint_background(slide, C_BG)
    _draw_slide_header(slide, slide_ir.get("title") or "文字起こし")

    blocks = slide_ir.get("blocks") or []
    if not blocks:
        _add_footer(slide)
        return

    block_h = Inches(0.9)
    y = Inches(1.5)
    for b in blocks:
        if y + block_h > Inches(SLIDE_H_IN - 0.6):
            break
        speaker = (b.get("speaker") or "").strip()
        time_label = (b.get("timeLabel") or "").strip()
        text = (b.get("text") or "").strip()
        if not text:
            continue

        header_parts = []
        if speaker:
            header_parts.append(f"🗣 {speaker}")
        if time_label:
            header_parts.append(f"⏱ {time_label}")
        header = "   ".join(header_parts)

        if header:
            _add_textbox(
                slide,
                left=MARGIN_X, top=y,
                width=Inches(SLIDE_W_IN) - MARGIN_X * 2, height=Inches(0.3),
                text=header,
                font_size=Pt(10), color=C_TEXT_MUTED, font=FONT_BODY,
            )
        _add_textbox(
            slide,
            left=MARGIN_X, top=y + Inches(0.3),
            width=Inches(SLIDE_W_IN) - MARGIN_X * 2, height=Inches(0.6),
            text=text,
            font_size=Pt(13), color=C_TEXT, font=FONT_BODY,
            word_wrap=True,
        )
        y += block_h

    _add_footer(slide)


# ---------------------------------------------------------------------------
# Shared drawing primitives
# ---------------------------------------------------------------------------

def _blank_slide(prs: Presentation):
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    return prs.slides.add_slide(blank)


def _paint_background(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(SLIDE_W_IN), Inches(SLIDE_H_IN),
    )
    _fill(bg, color)
    _no_line(bg)
    # Send-to-back via XML manipulation isn't strictly needed since it's added
    # first; every subsequent shape will paint on top.


def _draw_slide_header(slide, title: str) -> None:
    _add_textbox(
        slide,
        left=MARGIN_X, top=MARGIN_TOP,
        width=Inches(SLIDE_W_IN) - MARGIN_X * 2, height=Inches(0.6),
        text=title,
        font_size=Pt(22), bold=True, color=C_TEXT, font=FONT_HEADING,
    )
    # Accent underline
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_X, MARGIN_TOP + Inches(0.75),
        Inches(0.8), Inches(0.05),
    )
    _fill(rule, C_PRIMARY)
    _no_line(rule)


def _add_footer(slide, text: str | None = "Created by DeepNote") -> None:
    if not text:
        return
    _add_textbox(
        slide,
        left=MARGIN_X, top=Inches(SLIDE_H_IN - 0.45),
        width=Inches(SLIDE_W_IN) - MARGIN_X * 2, height=Inches(0.3),
        text=text,
        font_size=Pt(9), color=C_TEXT_MUTED, font=FONT_BODY,
        align=PP_ALIGN.RIGHT,
    )


def _draw_card_grid(slide, items: List[Dict[str, Any]], *, card_builder) -> None:
    """Lay out up to 6 cards in 2-col x 3-row grid."""
    if not items:
        _add_textbox(
            slide,
            left=MARGIN_X, top=Inches(3.2),
            width=Inches(SLIDE_W_IN) - MARGIN_X * 2, height=Inches(0.8),
            text="該当項目はありません",
            font_size=Pt(14), color=C_TEXT_MUTED, font=FONT_BODY,
            align=PP_ALIGN.CENTER,
        )
        return

    cols = 2
    rows = 3
    gap_x = Inches(0.3)
    gap_y = Inches(0.25)
    top = Inches(1.5)
    total_w = Inches(SLIDE_W_IN) - MARGIN_X * 2
    card_w = Emu((total_w - gap_x * (cols - 1)) // cols)
    total_h = Inches(SLIDE_H_IN - 2.4)
    card_h = Emu((total_h - gap_y * (rows - 1)) // rows)

    for idx, item in enumerate(items[: cols * rows]):
        r = idx // cols
        c = idx % cols
        x = MARGIN_X + (card_w + gap_x) * c
        y = top + (card_h + gap_y) * r
        card_builder(slide, x, y, card_w, card_h, item)


def _draw_card_base(slide, x: Emu, y: Emu, w: Emu, h: Emu) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _fill(shape, C_SURFACE)
    _set_line(shape, C_BORDER, width=Pt(0.5))


def _draw_card_accent_bar(slide, x: Emu, y: Emu, h: Emu, color: RGBColor) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    _fill(bar, color)
    _no_line(bar)


def _draw_chip_row(slide, *, left: Emu, top: Emu, width: Emu, chips: List[str]) -> None:
    x = left
    y = top
    for chip in chips:
        cw = _estimate_chip_width(chip, Pt(11))
        if x + cw > left + width:
            break
        _draw_chip(slide, x=x, y=y, width=cw, height=Inches(0.35),
                   text=chip, fg=C_TEXT, bg=C_CHIP_BG, font_size=Pt(11))
        x += cw + Inches(0.1)


def _draw_colored_chips(slide, *, left: Emu, top: Emu, width: Emu,
                         chips: List[Tuple[str, RGBColor]]) -> None:
    x = left
    for text, color in chips:
        cw = _estimate_chip_width(text, Pt(10))
        if x + cw > left + width:
            break
        _draw_chip(slide, x=x, y=top, width=cw, height=Inches(0.3),
                   text=text, fg=color, bg=C_CHIP_BG, font_size=Pt(10))
        x += cw + Inches(0.08)


def _draw_chip(slide, *, x: Emu, y: Emu, width: Emu, height: Emu,
               text: str, fg: RGBColor, bg: RGBColor, font_size) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    _fill(shape, bg)
    _no_line(shape)
    tf = shape.text_frame
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = font_size
    run.font.color.rgb = fg


def _estimate_chip_width(text: str, font_size) -> Emu:
    # Rough width estimation that's fine for layouts in the dark theme.
    size_pt = font_size.pt if hasattr(font_size, "pt") else 11
    per_char = Inches(size_pt * 0.018)
    # CJK-dense strings take more width.
    cjk_bias = sum(1 for ch in text if ord(ch) > 0x2E80) * 0.6
    w = Inches(0.3) + per_char * (len(text) + int(cjk_bias))
    return w


def _add_textbox(slide, *, left: Emu, top: Emu, width: Emu, height: Emu,
                 text: str, font_size, color: RGBColor, font: str,
                 bold: bool = False, italic: bool = False,
                 align=PP_ALIGN.LEFT, word_wrap: bool = False) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _fill(shape, color: RGBColor) -> None:
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_line(shape, color: RGBColor, *, width) -> None:
    line = shape.line
    line.color.rgb = color
    line.width = width


def _no_line(shape) -> None:
    line = shape.line
    line.fill.background()


# ---------------------------------------------------------------------------
# Label helpers
# ---------------------------------------------------------------------------

def _status_label(status: str) -> str:
    mapping = {
        "confirmed": "✅ 合意",
        "open": "🟠 保留",
        "rejected": "❌ 却下",
        "inferred": "💭 推定",
    }
    return mapping.get(status, status or "状態")


def _severity_label(severity: str) -> str:
    mapping = {
        "high": "HIGH",
        "medium": "MEDIUM",
        "low": "LOW",
    }
    return mapping.get(severity, severity.upper() if severity else "MEDIUM")
