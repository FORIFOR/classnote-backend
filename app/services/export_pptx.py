"""PPTX renderer for session export.

Generates a richly-styled slide deck that preserves every field of the summary
JSON. Layout:
- Slide 1: Title band + bottom line hero + outcome badge
- Slide 2: Highlights + participants (if any) + keywords
- Slide 3+: Content cards for each section (tables, cards, bullets), auto
  paginated so nothing overflows.
"""
import io
import logging
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)

# ─────────────────────────────────────────────────────
# Canvas + palette
# ─────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_X = Inches(0.6)
# Keep both EMU width and inch scalar. Arithmetic on python-pptx Length may
# yield raw int (EMU), while some layout logic needs float inches.
CONTENT_W_IN = 13.333 - 1.2
CONTENT_W = Inches(CONTENT_W_IN)

C_PRIMARY     = RGBColor(0x1F, 0x4F, 0xD8)
C_PRIMARY_SOFT = RGBColor(0xEE, 0xF3, 0xFF)
C_HEAD        = RGBColor(0x1A, 0x1A, 0x2E)
C_TEXT        = RGBColor(0x22, 0x22, 0x33)
C_CAPTION     = RGBColor(0x55, 0x55, 0x66)
C_MUTED       = RGBColor(0x88, 0x88, 0x99)
C_CARD_BG     = RGBColor(0xF7, 0xF8, 0xFB)
C_CARD_LINE   = RGBColor(0xE4, 0xE4, 0xEC)
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_BG          = RGBColor(0xFF, 0xFF, 0xFF)
C_SUCCESS     = RGBColor(0x16, 0xA3, 0x4A)
C_SUCCESS_BG  = RGBColor(0xEC, 0xFD, 0xF3)
C_WARN        = RGBColor(0xEA, 0x58, 0x0C)
C_WARN_BG     = RGBColor(0xFF, 0xF4, 0xED)
C_NEUTRAL     = RGBColor(0x6B, 0x72, 0x80)
C_NEUTRAL_BG  = RGBColor(0xF3, 0xF4, 0xF6)

EMPHASIS = {
    "primary": (C_PRIMARY, C_PRIMARY_SOFT),
    "success": (C_SUCCESS, C_SUCCESS_BG),
    "warning": (C_WARN, C_WARN_BG),
    "neutral": (C_NEUTRAL, C_NEUTRAL_BG),
}


# ─────────────────────────────────────────────────────
# Public entry
# ─────────────────────────────────────────────────────
def render_pptx(
    *,
    title: str,
    date_text: Optional[str] = None,
    duration_text: Optional[str] = None,
    mode: str = "meeting",
    highlights: Optional[List[Any]] = None,
    overview: Optional[str] = None,
    keywords: Optional[List[str]] = None,
    sections: Optional[List[Dict[str, Any]]] = None,
    transcript: Optional[str] = None,
    bottom_line: Optional[str] = None,
    why_it_matters: Optional[str] = None,
    outcome_status: Optional[str] = None,
    participants: Optional[List[Dict[str, str]]] = None,
    **_kwargs,
) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 2026-05-08 design refresh: PPTX is now a 1–2 slide compact deck
    # so it stays readable as a "summary attachment" rather than a
    # 10-slide per-section deck. Slide 1 is the title hero; slide 2
    # consolidates the highlights + key decisions/TODOs/overview
    # extracted from sections. The detailed per-section content stays
    # in the PDF / DOCX export.

    # Slide 1: Title + bottom line
    _slide_title(prs, title, date_text, duration_text, mode,
                 bottom_line, why_it_matters, outcome_status)

    # Slide 2: Compact "everything else" — only emit when we have
    # something to show.
    rich_hls = _normalize_highlights(highlights)
    digest_bullets = _condense_sections_to_bullets(sections or [], limit=8)
    if rich_hls or participants or overview or digest_bullets or keywords:
        _slide_compact_summary(
            prs,
            highlights=rich_hls,
            participants=participants,
            overview=overview,
            keywords=keywords,
            digest_bullets=digest_bullets,
        )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─────────────────────────────────────────────────────
# Compact summary helpers (2026-05-08 design refresh)
# ─────────────────────────────────────────────────────

def _condense_sections_to_bullets(
    sections: List[Dict[str, Any]],
    *,
    limit: int = 8,
) -> List[Tuple[str, str]]:
    """Flatten arbitrarily many sections into ``(label, text)`` pairs
    so the 2-slide deck can show a small snapshot. Stops at ``limit``
    bullets — anything past that lives in the PDF / DOCX exports.

    Recognised section variants (best-effort, robust to missing keys):
      - ``type: "bullets"`` → up to 3 from ``items``
      - ``type: "grouped_bullets"`` → flatten 1 group, up to 3 items
      - ``type: "table"`` → first column of up to 2 rows
      - ``type: "cards"`` → first ``title`` of up to 2 cards
      - ``type: "text"`` → first sentence
    """
    out: List[Tuple[str, str]] = []
    for sec in sections:
        label = (sec.get("title") or "").strip() or "—"
        st = (sec.get("type") or "").lower()
        items_added = 0
        if st in ("bullets", ""):
            for it in (sec.get("items") or [])[:3]:
                txt = it if isinstance(it, str) else (it.get("text") or "")
                if txt:
                    out.append((label, txt))
                    items_added += 1
        elif st == "grouped_bullets":
            groups = sec.get("groups") or []
            for g in groups[:1]:
                for it in (g.get("items") or [])[:3]:
                    txt = it if isinstance(it, str) else (it.get("text") or "")
                    if txt:
                        out.append((label, txt))
                        items_added += 1
        elif st == "table":
            for row in (sec.get("rows") or [])[:2]:
                if row:
                    out.append((label, str(row[0])))
                    items_added += 1
        elif st == "cards":
            for card in (sec.get("cards") or [])[:2]:
                t = (card.get("title") or "").strip()
                if t:
                    out.append((label, t))
                    items_added += 1
        elif st == "text":
            body = (sec.get("body") or "").strip()
            if body:
                first_line = body.splitlines()[0][:120]
                out.append((label, first_line))
        if len(out) >= limit:
            break
    return out[:limit]


def _slide_compact_summary(
    prs,
    *,
    highlights: List[Dict[str, Any]],
    participants: Optional[List[Dict[str, str]]],
    overview: Optional[str],
    keywords: Optional[List[str]],
    digest_bullets: List[Tuple[str, str]],
) -> None:
    """Single slide combining highlights + overview + digest bullets +
    participants + keywords. Items are deliberately compressed so a
    busy session still fits."""
    slide = _blank_slide(prs)

    # Header band
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.55))
    band.fill.solid()
    band.fill.fore_color.rgb = C_PRIMARY
    band.line.fill.background()
    head = slide.shapes.add_textbox(MARGIN_X, Inches(0.10), CONTENT_W, Inches(0.4))
    hp = head.text_frame.paragraphs[0]
    hp.text = "ハイライト & まとめ"
    hp.font.size = Pt(16)
    hp.font.bold = True
    hp.font.color.rgb = C_WHITE

    cur_y = Inches(0.85)
    inner_w = CONTENT_W

    # Overview (first paragraph only)
    if overview:
        first = (overview or "").strip().splitlines()[0][:200]
        if first:
            ov = slide.shapes.add_textbox(MARGIN_X, cur_y, inner_w, Inches(0.65))
            ovp = ov.text_frame.paragraphs[0]
            ovp.text = first
            ovp.font.size = Pt(12)
            ovp.font.color.rgb = C_TEXT
            cur_y += Inches(0.7)

    # Two columns: highlights (left) + digest bullets (right)
    col_w = (CONTENT_W - Inches(0.3)) / 2
    left_x = MARGIN_X
    right_x = MARGIN_X + col_w + Inches(0.3)
    list_h = Inches(3.4)

    # Left column: highlights
    if highlights:
        lt = slide.shapes.add_textbox(left_x, cur_y, col_w, Inches(0.32))
        ltp = lt.text_frame.paragraphs[0]
        ltp.text = "Highlights"
        ltp.font.size = Pt(11)
        ltp.font.bold = True
        ltp.font.color.rgb = C_PRIMARY
        body = slide.shapes.add_textbox(left_x, cur_y + Inches(0.35), col_w, list_h)
        bf = body.text_frame
        bf.word_wrap = True
        for i, h in enumerate(highlights[:5]):
            txt = h.get("text") if isinstance(h, dict) else str(h)
            if not txt: continue
            para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            para.text = f"• {txt[:80]}"
            para.font.size = Pt(11)
            para.font.color.rgb = C_TEXT

    # Right column: digest bullets (decisions / todos / etc.)
    if digest_bullets:
        rt = slide.shapes.add_textbox(right_x, cur_y, col_w, Inches(0.32))
        rtp = rt.text_frame.paragraphs[0]
        rtp.text = "決定事項・TODO・メモ"
        rtp.font.size = Pt(11)
        rtp.font.bold = True
        rtp.font.color.rgb = C_PRIMARY
        body = slide.shapes.add_textbox(right_x, cur_y + Inches(0.35), col_w, list_h)
        bf = body.text_frame
        bf.word_wrap = True
        seen_label = ""
        idx = 0
        for label, txt in digest_bullets[:6]:
            para = bf.paragraphs[0] if idx == 0 else bf.add_paragraph()
            prefix = f"[{label[:8]}] " if label and label != seen_label else "• "
            seen_label = label
            para.text = f"{prefix}{txt[:80]}"
            para.font.size = Pt(11)
            para.font.color.rgb = C_TEXT
            idx += 1

    # Footer: participants + keywords on the same row
    foot_y = SLIDE_H - Inches(0.85)
    if participants:
        names = [p.get("name") or "(不明)" for p in participants[:6]]
        ptxt = "参加者: " + " / ".join(names)
        if len(participants) > 6:
            ptxt += f" 他 {len(participants)-6} 名"
        pbox = slide.shapes.add_textbox(MARGIN_X, foot_y, CONTENT_W, Inches(0.3))
        pp = pbox.text_frame.paragraphs[0]
        pp.text = ptxt
        pp.font.size = Pt(9)
        pp.font.color.rgb = C_MUTED
        foot_y += Inches(0.25)
    if keywords:
        kbox = slide.shapes.add_textbox(MARGIN_X, foot_y, CONTENT_W, Inches(0.3))
        kp = kbox.text_frame.paragraphs[0]
        kp.text = "Keywords: " + " · ".join(keywords[:8])
        kp.font.size = Pt(9)
        kp.font.color.rgb = C_PRIMARY


# ─────────────────────────────────────────────────────
# Slide 1 — Title hero
# ─────────────────────────────────────────────────────
def _slide_title(prs, title, date_text, duration_text, mode,
                 bottom_line, why_it_matters, outcome_status):
    slide = _blank_slide(prs)

    # Dark band with gradient effect (dark navy)
    band_h = Inches(3.2)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, band_h)
    band.fill.solid()
    band.fill.fore_color.rgb = C_HEAD
    band.line.fill.background()

    # Accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, band_h, SLIDE_W, Inches(0.08))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = C_PRIMARY
    stripe.line.fill.background()

    # Mode label
    mode_label = _mode_label(mode).upper()
    lbl = slide.shapes.add_textbox(MARGIN_X, Inches(0.5), Inches(3), Inches(0.3))
    lp = lbl.text_frame.paragraphs[0]
    lp.text = mode_label
    lp.font.size = Pt(11)
    lp.font.bold = True
    lp.font.color.rgb = C_PRIMARY

    # Title
    ttxb = slide.shapes.add_textbox(MARGIN_X, Inches(0.9), CONTENT_W, Inches(1.3))
    ttf = ttxb.text_frame
    ttf.word_wrap = True
    p = ttf.paragraphs[0]
    p.text = title or "Untitled"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = C_WHITE

    # Meta line
    meta_parts = [x for x in [date_text, duration_text] if x]
    if meta_parts:
        mtxb = slide.shapes.add_textbox(MARGIN_X, Inches(2.35), CONTENT_W, Inches(0.35))
        mp = mtxb.text_frame.paragraphs[0]
        mp.text = " · ".join(meta_parts)
        mp.font.size = Pt(13)
        mp.font.color.rgb = RGBColor(0xBB, 0xBB, 0xCC)

    # Hero callout below band
    if bottom_line or why_it_matters or outcome_status:
        callout_top = Inches(3.7)
        callout = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_X, callout_top,
            CONTENT_W, Inches(3.0),
        )
        callout.fill.solid()
        callout.fill.fore_color.rgb = C_PRIMARY_SOFT
        callout.line.color.rgb = C_PRIMARY
        callout.line.width = Pt(1.5)

        inner_x = MARGIN_X + Inches(0.4)
        inner_w = CONTENT_W - Inches(0.8)
        cur_y = callout_top + Inches(0.35)

        if outcome_status:
            badge_label = _outcome_label(outcome_status)
            btxb = slide.shapes.add_textbox(inner_x, cur_y, inner_w, Inches(0.3))
            bp = btxb.text_frame.paragraphs[0]
            bp.text = badge_label
            bp.font.size = Pt(12)
            bp.font.bold = True
            bp.font.color.rgb = C_PRIMARY
            cur_y += Inches(0.4)

        if bottom_line:
            btxb = slide.shapes.add_textbox(inner_x, cur_y, inner_w, Inches(1.5))
            tf = btxb.text_frame
            tf.word_wrap = True
            bp = tf.paragraphs[0]
            bp.text = bottom_line
            bp.font.size = Pt(22)
            bp.font.bold = True
            bp.font.color.rgb = C_HEAD
            cur_y += Inches(1.5)

        if why_it_matters:
            wtxb = slide.shapes.add_textbox(inner_x, cur_y, inner_w, Inches(0.8))
            tf = wtxb.text_frame
            tf.word_wrap = True
            wp = tf.paragraphs[0]
            wp.text = why_it_matters
            wp.font.size = Pt(13)
            wp.font.color.rgb = C_CAPTION

    _add_footer(slide, page_label=None)


# ─────────────────────────────────────────────────────
# Slide 2 — Highlights + participants
# ─────────────────────────────────────────────────────
def _slide_highlights(prs, highlights, participants, overview, keywords):
    slide = _blank_slide(prs)
    _slide_header(slide, "要点 / 概要")

    body_top = Inches(1.1)
    body_h = Inches(5.3)

    has_participants = bool(participants)
    left_w = CONTENT_W * (0.62 if has_participants else 1.0)

    # Left column — highlights
    if highlights:
        _small_label(slide, MARGIN_X, body_top, "要点", C_PRIMARY)
        htxb = slide.shapes.add_textbox(
            MARGIN_X, body_top + Inches(0.35),
            left_w, body_h - Inches(0.35),
        )
        tf = htxb.text_frame
        tf.word_wrap = True
        for i, h in enumerate(highlights[:10]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            category = h.get("category")
            prefix = f"[{_category_label(category)}] " if category else ""
            text = (h.get("text") or "")
            p.text = f"•  {prefix}{text}"
            p.font.size = Pt(13)
            p.font.color.rgb = C_TEXT
            p.space_after = Pt(6)

    # Right column — participants
    if has_participants:
        right_x = MARGIN_X + left_w + Inches(0.4)
        right_w = CONTENT_W - left_w - Inches(0.4)
        _small_label(slide, right_x, body_top, "参加者", C_PRIMARY)
        pt_y = body_top + Inches(0.35)
        max_rows = 10
        row_h = Inches(0.42)
        for i, person in enumerate(participants[:max_rows]):
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, right_x, pt_y + row_h * i,
                right_w, Inches(0.38),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = C_CARD_BG
            card.line.color.rgb = C_CARD_LINE
            card.line.width = Pt(0.5)
            txb = slide.shapes.add_textbox(
                right_x + Inches(0.15), pt_y + row_h * i + Inches(0.05),
                right_w - Inches(0.3), Inches(0.3),
            )
            tf = txb.text_frame
            p = tf.paragraphs[0]
            name_run = p.add_run()
            name_run.text = person.get("name") or "(不明)"
            name_run.font.size = Pt(11)
            name_run.font.bold = True
            name_run.font.color.rgb = C_HEAD
            if person.get("role"):
                role_run = p.add_run()
                role_run.text = f"  {person['role']}"
                role_run.font.size = Pt(10)
                role_run.font.color.rgb = C_MUTED

    # Overview as single paragraph under highlights if present & no participants column
    if overview and not highlights:
        otxb = slide.shapes.add_textbox(MARGIN_X, body_top + Inches(0.3), CONTENT_W, Inches(4))
        tf = otxb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = overview
        p.font.size = Pt(14)
        p.font.color.rgb = C_TEXT

    # Keywords chip strip at bottom
    if keywords:
        kw_y = Inches(6.55)
        kw_label = slide.shapes.add_textbox(MARGIN_X, kw_y, Inches(1.5), Inches(0.25))
        kl = kw_label.text_frame.paragraphs[0]
        kl.text = "キーワード"
        kl.font.size = Pt(9)
        kl.font.bold = True
        kl.font.color.rgb = C_MUTED

        kw_txb = slide.shapes.add_textbox(
            MARGIN_X + Inches(1.4), kw_y,
            CONTENT_W - Inches(1.4), Inches(0.3),
        )
        kp = kw_txb.text_frame.paragraphs[0]
        kp.text = "   ".join(f"#{k}" for k in keywords[:14])
        kp.font.size = Pt(10)
        kp.font.color.rgb = C_PRIMARY

    _add_footer(slide)


# ─────────────────────────────────────────────────────
# Section slides
# ─────────────────────────────────────────────────────
def _render_section_slide(prs, sec: Dict[str, Any]) -> None:
    sec_type = sec.get("type", "text")
    sec_title = sec.get("title", "")
    emphasis = sec.get("emphasis", "primary")
    accent, accent_bg = EMPHASIS.get(emphasis, EMPHASIS["primary"])

    if sec_type == "text":
        _text_slide(prs, sec_title, sec.get("body", ""), accent)

    elif sec_type == "bullets":
        items = sec.get("items", [])
        _paginated_bullets(prs, sec_title, [(None, items)], accent, accent_bg)

    elif sec_type == "grouped_bullets":
        groups = [(g.get("label"), g.get("items", [])) for g in sec.get("groups", [])]
        _paginated_bullets(prs, sec_title, groups, accent, accent_bg)

    elif sec_type == "table":
        columns = sec.get("columns", [])
        rows = sec.get("rows", [])
        if columns and rows:
            _table_slide(prs, sec_title, columns, rows, accent, accent_bg)

    elif sec_type == "cards":
        cards = sec.get("cards", [])
        if cards:
            _cards_slide(prs, sec_title, cards, accent, accent_bg)


def _text_slide(prs, title, body, accent):
    slide = _blank_slide(prs)
    _slide_header(slide, title, accent=accent)
    txb = slide.shapes.add_textbox(MARGIN_X, Inches(1.1), CONTENT_W, Inches(5.5))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for para in str(body).split("\n"):
        para = para.strip()
        if not para:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = para
        p.font.size = Pt(14)
        p.font.color.rgb = C_TEXT
        p.space_after = Pt(8)
    _add_footer(slide)


def _paginated_bullets(prs, title, groups: List[Tuple[Optional[str], List[str]]],
                        accent: RGBColor, accent_bg: RGBColor) -> None:
    """Render bullets (possibly multi-group) split into multiple slides if long."""
    MAX_LINES_PER_SLIDE = 18
    # Flatten groups into lines with optional group header markers
    lines: List[Tuple[str, str]] = []  # (kind, text); kind in {"label", "item"}
    for label, items in groups:
        if label:
            lines.append(("label", label))
        for it in items:
            lines.append(("item", str(it)))

    # Split into slide-sized chunks, keeping labels with following items
    slides_lines: List[List[Tuple[str, str]]] = []
    current: List[Tuple[str, str]] = []
    count_items = 0
    for kind, text in lines:
        if kind == "item":
            if count_items >= MAX_LINES_PER_SLIDE:
                slides_lines.append(current)
                current = []
                count_items = 0
            current.append((kind, text))
            count_items += 1
        else:  # label
            current.append((kind, text))
    if current:
        slides_lines.append(current)

    total = len(slides_lines)
    for idx, chunk in enumerate(slides_lines):
        slide = _blank_slide(prs)
        heading = title if total == 1 else f"{title}  ({idx + 1}/{total})"
        _slide_header(slide, heading, accent=accent)
        txb = slide.shapes.add_textbox(MARGIN_X, Inches(1.1), CONTENT_W, Inches(5.5))
        tf = txb.text_frame
        tf.word_wrap = True
        first = True
        for kind, text in chunk:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            if kind == "label":
                p.text = text
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = accent
                p.space_before = Pt(6)
                p.space_after = Pt(3)
            else:
                p.text = f"•  {text}"
                p.font.size = Pt(13)
                p.font.color.rgb = C_TEXT
                p.space_after = Pt(5)
        _add_footer(slide)


def _table_slide(prs, title, columns: List[str], rows: List[list],
                 accent: RGBColor, accent_bg: RGBColor) -> None:
    MAX_ROWS_PER_SLIDE = 10
    total_chunks = (len(rows) + MAX_ROWS_PER_SLIDE - 1) // MAX_ROWS_PER_SLIDE
    for idx in range(total_chunks):
        chunk = rows[idx * MAX_ROWS_PER_SLIDE : (idx + 1) * MAX_ROWS_PER_SLIDE]
        slide = _blank_slide(prs)
        heading = title if total_chunks == 1 else f"{title}  ({idx + 1}/{total_chunks})"
        _slide_header(slide, heading, accent=accent)

        top = Inches(1.2)
        height = Inches(5.5)
        cols = len(columns)
        rows_count = len(chunk) + 1  # +1 for header
        tbl_shape = slide.shapes.add_table(rows_count, cols, MARGIN_X, top, CONTENT_W, height)
        tbl = tbl_shape.table

        # Distribute column widths — first column wider
        if cols == 2:
            widths = [0.35, 0.65]
        elif cols == 3:
            widths = [0.3, 0.45, 0.25]
        elif cols == 4:
            widths = [0.3, 0.25, 0.25, 0.2]
        elif cols == 5:
            widths = [0.3, 0.22, 0.18, 0.15, 0.15]
        else:
            widths = [1.0 / cols] * cols
        for i, w in enumerate(widths[:cols]):
            tbl.columns[i].width = Inches(float(CONTENT_W.inches) * w)

        # Header
        for i, col in enumerate(columns):
            cell = tbl.cell(0, i)
            cell.text = col
            cell.fill.solid()
            cell.fill.fore_color.rgb = accent_bg
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(11)
                    run.font.bold = True
                    run.font.color.rgb = accent

        # Body
        for r_idx, row_data in enumerate(chunk, start=1):
            padded = list(row_data) + [""] * (cols - len(row_data))
            for c_idx, cell_text in enumerate(padded[:cols]):
                cell = tbl.cell(r_idx, c_idx)
                cell.text = str(cell_text or "")
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE if r_idx % 2 == 1 else C_CARD_BG
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(10)
                        run.font.color.rgb = C_TEXT

        _add_footer(slide)


def _cards_slide(prs, title, cards: List[Dict[str, Any]],
                 accent: RGBColor, accent_bg: RGBColor) -> None:
    MAX_CARDS_PER_SLIDE = 4
    total_chunks = (len(cards) + MAX_CARDS_PER_SLIDE - 1) // MAX_CARDS_PER_SLIDE
    for idx in range(total_chunks):
        chunk = cards[idx * MAX_CARDS_PER_SLIDE : (idx + 1) * MAX_CARDS_PER_SLIDE]
        slide = _blank_slide(prs)
        heading = title if total_chunks == 1 else f"{title}  ({idx + 1}/{total_chunks})"
        _slide_header(slide, heading, accent=accent)

        n = len(chunk)
        cols = 2 if n > 1 else 1
        rows_count = (n + cols - 1) // cols
        gap = Inches(0.25)
        body_top = Inches(1.15)
        body_h = Inches(5.6)
        card_w = (CONTENT_W - gap * (cols - 1)) / cols
        card_h = (body_h - gap * (rows_count - 1)) / rows_count

        for i, card in enumerate(chunk):
            col = i % cols
            row = i // cols
            x = MARGIN_X + col * (card_w + gap)
            y = body_top + row * (card_h + gap)
            _draw_card(slide, x, y, card_w, card_h, card, accent)

        _add_footer(slide)


def _draw_card(slide, x, y, w, h, card: Dict[str, Any], accent: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_CARD_BG
    shape.line.color.rgb = C_CARD_LINE
    shape.line.width = Pt(0.75)

    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    inner_pad = Inches(0.2)
    content_w = w - inner_pad * 2
    cur_y = y + Inches(0.2)

    # Card title
    title = card.get("title", "")
    if title:
        ttxb = slide.shapes.add_textbox(x + inner_pad, cur_y, content_w, Inches(0.5))
        tf = ttxb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_HEAD
        cur_y += Inches(0.5)

    # Fields
    fields = [(k, v) for k, v in card.get("fields", []) if v]
    if not fields:
        return
    field_h = h - (cur_y - y) - Inches(0.2)
    ftxb = slide.shapes.add_textbox(x + inner_pad, cur_y, content_w, field_h)
    tf = ftxb.text_frame
    tf.word_wrap = True
    first = True
    for key, value in fields:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if key:
            kr = p.add_run()
            kr.text = f"{key}  "
            kr.font.size = Pt(9)
            kr.font.color.rgb = C_MUTED
            vr = p.add_run()
            vr.text = str(value)
            vr.font.size = Pt(11)
            vr.font.color.rgb = C_TEXT
        else:
            p.text = str(value)
            p.font.size = Pt(11)
            p.font.color.rgb = C_TEXT
        p.space_after = Pt(3)


# ─────────────────────────────────────────────────────
# Slide primitives
# ─────────────────────────────────────────────────────
def _blank_slide(prs):
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C_BG
    return slide


def _slide_header(slide, title: str, accent: RGBColor = C_PRIMARY) -> None:
    # Accent stripe on the left
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_X, Inches(0.35), Inches(0.12), Inches(0.5),
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()

    txb = slide.shapes.add_textbox(
        MARGIN_X + Inches(0.25), Inches(0.3),
        CONTENT_W - Inches(0.25), Inches(0.6),
    )
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C_HEAD

    # Bottom hairline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_X, Inches(0.95), CONTENT_W, Pt(1),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C_CARD_LINE
    line.line.fill.background()


def _small_label(slide, x, y, text: str, color: RGBColor) -> None:
    txb = slide.shapes.add_textbox(x, y, Inches(3), Inches(0.3))
    p = txb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = color


def _add_footer(slide, page_label: Optional[str] = "Created by DeepNote") -> None:
    if not page_label:
        return
    txb = slide.shapes.add_textbox(MARGIN_X, Inches(7.15), CONTENT_W, Inches(0.25))
    p = txb.text_frame.paragraphs[0]
    p.text = page_label
    p.font.size = Pt(8)
    p.font.color.rgb = C_MUTED
    p.alignment = PP_ALIGN.RIGHT


# ─────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────
def _normalize_highlights(highlights) -> List[Dict[str, Any]]:
    if not highlights:
        return []
    out = []
    for h in highlights:
        if isinstance(h, str):
            out.append({"text": h, "category": None, "needConfirm": False})
        elif isinstance(h, dict):
            out.append({
                "text": h.get("text") or h.get("title") or "",
                "category": h.get("category"),
                "needConfirm": bool(h.get("needConfirm")),
            })
    return [x for x in out if x.get("text")]


def _mode_label(mode: str) -> str:
    return {"lecture": "講義", "meeting": "会議", "translate": "翻訳"}.get(mode, mode)


def _category_label(cat: str) -> str:
    return {
        "decision": "決定",
        "todo": "TODO",
        "concern": "懸念",
        "insight": "気づき",
        "fact": "事実",
        "info": "情報",
        "risk": "リスク",
    }.get(cat, cat)


def _outcome_label(status: str) -> str:
    return {
        "action_agreed": "✓ 合意・次アクション確定",
        "tentative_alignment": "△ 暫定合意",
        "discussion_only": "・ 議論のみ",
        "unresolved": "! 未解決",
    }.get(status, status)
